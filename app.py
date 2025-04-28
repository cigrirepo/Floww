# app.py
# Floww – AI-powered Deal Workflow & Proposal Generator
# ====================================================
#  • Real-company enrichment (ticker / Clearbit)
#  • Two-pane tabs (Workflow | Proposal)
#  • Advanced Mermaid diagrams with theming + JSON editor
#  • Competitor benchmarks, CRM playbook, PPTX export
#  • Spreadsheet-style pricing grid, presets, totals
#  • AI-generated proposal → PDF export
#  • Robust handling of pricing formats in AI output

import io, os, re
import json
from datetime import date
from typing import List, Dict, Any

import pandas as pd
import streamlit as st
from pydantic import BaseModel, ValidationError
from openai import OpenAI
import yfinance as yf
import requests
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from pptx import Presentation

# ── Config ────────────────────────────────────────────
st.set_page_config(page_title="Floww Advanced", layout="wide")
st.title("Floww Advanced")
st.caption("AI-driven workflows, diagrams & proposals tailored to real companies")

# ── API Clients ───────────────────────────────────────
api_key      = os.getenv("OPENAI_API_KEY")
clearbit_key = os.getenv("CLEARBIT_KEY")
if not api_key:
    st.error("Missing OPENAI_API_KEY")
    st.stop()
client = OpenAI(api_key=api_key)

# ── Data Models ───────────────────────────────────────
class StageModel(BaseModel):
    stage: str
    tip:   str

class WorkflowModel(BaseModel):
    workflow: List[StageModel]

class ProposalModel(BaseModel):
    title: str
    executive_summary: str
    background: str
    solution_overview: str
    deliverables: List[str]
    pricing: Any
    next_steps: str

# ── Sidebar – Company info & presets ──────────────────
st.sidebar.header("Company Info")
company   = st.sidebar.text_input("Company name", "Acme Corp")
website   = st.sidebar.text_input("Website URL (optional)")
ticker    = st.sidebar.text_input("Stock ticker (optional)")
industry  = st.sidebar.selectbox("Industry", ["Fintech","SaaS","Retail","Healthcare","Other"])
persona   = st.sidebar.selectbox("Persona", ["Enterprise AE","SMB SDR","Partner Manager"])

# Fetch description / market cap when possible
description = {
    "Fintech":"a regulated financial services provider",
    "SaaS":"a fast-growing SaaS company",
    "Retail":"a global retail chain",
    "Healthcare":"a healthcare technology provider",
    "Other":"a leading enterprise"
}[industry]
market_cap = None
if ticker:
    try:
        info = yf.Ticker(ticker).info
        description = info.get("longBusinessSummary", description)
        cap = info.get("marketCap")
        market_cap = f"${cap:,.0f}" if cap else None
    except:
        pass
elif website and clearbit_key:
    try:
        r = requests.get(
            f"https://company.clearbit.com/v1/domains/find?domain={website}",
            headers={"Authorization": f"Bearer {clearbit_key}"}, timeout=4
        )
        description = r.json().get("description", description)
    except:
        pass

st.sidebar.header("Advanced Options")
crm_file   = st.sidebar.file_uploader("CRM CSV (lead,stage,probability)", type="csv")
competitor = st.sidebar.selectbox("Competitor Benchmark", ["None","Visa","Stripe","Amex"])
pptx_ok    = st.sidebar.checkbox("Enable PPTX Export", True)

# ── Tabs ──────────────────────────────────────────────
tab_wf, tab_prop = st.tabs(["Workflow","Proposal"])

# -----------------------------------------------------#
#                   WORKFLOW TAB                       #
# -----------------------------------------------------#
with tab_wf:
    if st.button("Generate Workflow"):
        system = (
            "You are an enterprise sales consultant. "
            "Return ONLY valid JSON: {\"workflow\":[{\"stage\":\"\",\"tip\":\"\"}]}"
        )
        prompt = (
            f"Build a deal-closing workflow for a {persona} at {company}. "
            f"Company description: {description[:200]}. "
            f"Industry: {industry}. Market cap: {market_cap or 'N/A'}."
        )
        resp = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role":"system","content":system},
                      {"role":"user","content":prompt}],
            temperature=0.0,
            max_tokens=600
        )
        raw = resp.choices[0].message.content
        match = re.search(r"\{.*\}", raw, re.S)
        if not match:
            st.error("AI did not return JSON")
            st.code(raw)
            st.stop()
        wf_json = match.group(0)
        try:
            wf = WorkflowModel.parse_raw(wf_json)
        except ValidationError as e:
            st.error("Invalid JSON")
            st.code(wf_json)
            st.error(str(e))
            st.stop()
        st.session_state['wf_json']  = wf_json
        st.session_state['stages']   = [s.stage for s in wf.workflow]
        st.session_state['tips']     = {s.stage: s.tip for s in wf.workflow}

    if 'wf_json' in st.session_state:
        st.subheader("Edit workflow JSON (live)")
        new_json = st.text_area("JSON", st.session_state['wf_json'], height=200)
        if st.button("Re-Render Diagram"):
            try:
                wf = WorkflowModel.parse_raw(new_json)
                st.session_state['wf_json']  = new_json
                st.session_state['stages']   = [s.stage for s in wf.workflow]
                st.session_state['tips']     = {s.stage: s.tip for s in wf.workflow}
            except:
                st.error("Invalid JSON")

    if 'stages' in st.session_state:
        theme = (
            "%%{init:{'themeVariables':{'primaryColor':'#00ADEF',"
            "'lineColor':'#888'}}}%%\n"
        )
        sys_d = "You are a diagram expert. Return Mermaid flowchart only."
        usr_d = f"Stages: {st.session_state['stages']}. Company: {company}."
        m_resp = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role":"system","content":sys_d},
                      {"role":"user","content":usr_d}],
            temperature=0.0,
            max_tokens=200
        )
        mcode = re.sub(r"^```(?:mermaid)?|```$", "", m_resp.choices[0].message.content, flags=re.M).strip()
        full_code = theme + mcode
        st.markdown("##### Mermaid Diagram")
        st.code(full_code)
        st.markdown(f"```mermaid\n{full_code}\n```", unsafe_allow_html=True)

    if 'stages' in st.session_state and competitor != 'None':
        bench = {
            'Visa': ['Prospecting','KYC Check','Regulatory Review'],
            'Stripe': ['API Integration Pilot','Sandbox Testing'],
            'Amex': ['Regulatory Compliance','Fraud Assessment']
        }
        st.markdown(f"#### Benchmark vs {competitor}")
        st.write(bench.get(competitor, []))

    if 'stages' in st.session_state and crm_file:
        df = pd.read_csv(crm_file)
        play = []
        for _, row in df.iterrows():
            tip_q = f"Lead data: {row.to_dict()}. Suggest next step."
            tip_resp = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role":"user","content":tip_q}],
                temperature=0.7,
                max_tokens=100
            )
            play.append({**row.to_dict(), 'suggestion': tip_resp.choices[0].message.content})
        st.markdown("#### Personalized Playbook from CRM")
        st.dataframe(pd.DataFrame(play))

    if 'stages' in st.session_state and pptx_ok:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = f"Floww Playbook: {company}"
        for s in st.session_state['stages']:
            sl = prs.slides.add_slide(prs.slide_layouts[1])
            sl.shapes.title.text = s
            sl.placeholders[1].text = st.session_state['tips'][s]
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        st.download_button(
            "Download PPTX Playbook",
            buf,
            "floww_playbook.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

# -----------------------------------------------------#
#                      PROPOSAL TAB                    #
# -----------------------------------------------------#
with tab_prop:
    st.subheader("Generate Proposal")

    col1, col2 = st.columns(2)
    with col1:
        client_name = st.text_input("Proposal for (client)", key="client")
    with col2:
        prop_date = st.date_input("Date", date.today())

    deliverables_txt = st.text_area(
        "Key deliverables (one per line)",
        placeholder="Integration setup\n24/7 support\nTraining sessions",
        height=100
    )

    st.markdown("#### Pricing Table")
    if "price_table" not in st.session_state:
        st.session_state["price_table"] = pd.DataFrame(
            [{"Item":"Integration","Qty":1,"Unit":"Lot","Unit Price":10000}]
        )

    p1, p2, p3 = st.columns(3)
    with p1:
        if st.button("🌱 Starter"):
            st.session_state["price_table"] = pd.DataFrame([
                {"Item":"Discovery","Qty":1,"Unit":"Lot","Unit Price":5000},
                {"Item":"Integration","Qty":1,"Unit":"Lot","Unit Price":20000},
            ])
    with p2:
        if st.button("🚀 Growth"):
            st.session_state["price_table"] = pd.DataFrame([
                {"Item":"Discovery","Qty":1,"Unit":"Lot","Unit Price":10000},
                {"Item":"Integration","Qty":1,"Unit":"Lot","Unit Price":50000},
                {"Item":"Training","Qty":1,"Unit":"Lot","Unit Price":15000},
            ])
    with p3:
        if st.button("🏢 Enterprise"):
            st.session_state["price_table"] = pd.DataFrame([
                {"Item":"Enterprise License (12 mo)","Qty":1,"Unit":"Lot","Unit Price":120000},
                {"Item":"Dedicated CSM","Qty":1,"Unit":"Yr","Unit Price":30000},
            ])

    price_df = st.data_editor(
        st.session_state["price_table"],
        num_rows="dynamic",
        use_container_width=True,
        key="price_editor"
    )
    price_df["Subtotal"] = price_df["Qty"] * price_df["Unit Price"]
    total = price_df["Subtotal"].sum()
    st.metric("Grand Total", f"${total:,.0f}")

    if st.button("Generate Proposal", type="primary"):
        # clean pricing to native Python types
        price_clean = price_df[["Item","Qty","Unit","Unit Price"]].astype({
            "Qty": int,
            "Unit Price": float
        })
        records = price_clean.to_dict(orient="records")

        # assemble plain-text prompt
        lines = [
            f"Client: {client_name}",
            f"Company: {company}",
            f"Date: {prop_date}",
            "",
            "Deliverables:"
        ] + [f"- {d}" for d in deliverables_txt.splitlines()] + [
            "",
            "Pricing:"
        ]
        for r in records:
            lines.append(f"- {r['Item']}, Qty: {r['Qty']}, Unit: {r['Unit']}, Price: ${r['Unit Price']:,.0f}")
        lines.append(f"\nTotal: ${total:,.0f}")

        sys_p = (
            "You are an expert sales engineer. "
            "Return ONLY valid JSON with keys: "
            "title, executive_summary, background, solution_overview, "
            "deliverables (list of strings), "
            "pricing (list of objects with keys item, qty, unit, price), "
            "next_steps."
        )
        user_p = "\n".join(lines)

        with st.spinner("Crafting proposal with Floww AI …"):
            resp = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role":"system","content":sys_p},
                          {"role":"user","content":user_p}],
                temperature=0.2,
                max_tokens=750,
            )
        prop_raw = resp.choices[0].message.content

        try:
            proposal = ProposalModel.parse_raw(prop_raw)
        except ValidationError as e:
            st.error("AI returned invalid JSON")
            st.code(prop_raw, language="json")
            st.error(str(e))
            st.stop()

        # normalize if pricing came as dict
        if isinstance(proposal.pricing, dict):
            normalized = []
            for item, details in proposal.pricing.items():
                if isinstance(details, dict):
                    normalized.append({
                        "item": item,
                        "qty": details.get("qty") or details.get("Qty") or 1,
                        "unit": details.get("unit") or details.get("Unit") or "",
                        "price": details.get("price") or details.get("Unit Price") or 0.0
                    })
            proposal.pricing = normalized

        st.success("Proposal generated ✔︎")
        st.markdown(f"## {proposal.title}")
        st.markdown(f"### Executive Summary\n{proposal.executive_summary}")
        st.markdown(f"### Background\n{proposal.background}")
        st.markdown(f"### Solution Overview\n{proposal.solution_overview}")
        st.markdown("### Deliverables")
        st.markdown("\n".join(f"- {d}" for d in proposal.deliverables))
        st.markdown("### Pricing")
        display_df = pd.DataFrame(proposal.pricing)
        st.table(display_df)
        st.markdown(f"### Next Steps\n{proposal.next_steps}")

        # PDF export
        buf = io.BytesIO()
        c   = canvas.Canvas(buf, pagesize=letter)
        w, h = letter
        y = h - 40
        def add_text(text, y_pos, indent=40, leading=14):
            for line in text.split("\n"):
                c.drawString(indent, y_pos, line)
                y_pos -= leading
                if y_pos < 60:
                    c.showPage()
                    y_pos = h - 40
            return y_pos

        c.setFont("Helvetica-Bold", 16)
        c.drawString(40, y, proposal.title)
        y -= 25
        c.setFont("Helvetica", 12)

        for sec in ["executive_summary","background","solution_overview","next_steps"]:
            c.setFont("Helvetica-Bold", 12)
            y = add_text(sec.replace("_", " ").title() + ":", y, 40)
            y -= 5
            c.setFont("Helvetica", 12)
            y = add_text(getattr(proposal, sec), y, 50)
            y -= 10

        c.setFont("Helvetica-Bold", 12)
        y = add_text("Pricing:", y, 40)
        y -= 5
        for row in display_df.itertuples(index=False):
            line = f"{row.item} ×{row.qty} {row.unit} … ${row.price:,.0f}"
            y = add_text(line, y, 50)
        add_text(f"Grand Total: ${total:,.0f}", y, 50)

        c.save()
        buf.seek(0)
        st.download_button(
            "Download Proposal PDF",
            buf,
            "proposal.pdf",
            "application/pdf"
        )

# ── Footer ───────────────────────────────────────────
st.sidebar.markdown("---")
st.sidebar.markdown("Powered by Adam Cigri & OpenAI")
